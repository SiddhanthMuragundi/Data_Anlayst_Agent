from flask import Flask, request, jsonify
import asyncio
import json
import base64
import io
from io import StringIO, BytesIO
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import pandas as pd
import numpy as np
import requests
from datetime import datetime
import re
import traceback
import seaborn as sns
from typing import Dict, List, Any, Optional, Union
from anthropic import Anthropic
import os
from bs4 import BeautifulSoup
import logging
import time
import warnings
import sqlite3
import csv
from urllib.parse import urlparse, urljoin
import zipfile
import tempfile
from dotenv import load_dotenv
import threading
from functools import partial
import unicodedata
import chardet
import mimetypes
import signal
from contextlib import contextmanager
from concurrent.futures import ThreadPoolExecutor, TimeoutError as FuturesTimeoutError

# Optional imports - handle gracefully if not installed
try:
    from playwright.async_api import async_playwright
    PLAYWRIGHT_AVAILABLE = True
except ImportError:
    PLAYWRIGHT_AVAILABLE = False
    async_playwright = None

try:
    import aiohttp
    AIOHTTP_AVAILABLE = True
except ImportError:
    AIOHTTP_AVAILABLE = False
    aiohttp = None

try:
    import xlrd
except ImportError:
    xlrd = None

try:
    from docx import Document
except ImportError:
    Document = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    import easyocr
    import cv2
    EASYOCR_AVAILABLE = True
except ImportError:
    EASYOCR_AVAILABLE = False
    easyocr = None
    cv2 = None

try:
    import fitz  # PyMuPDF for PDF processing
    PYMUPDF_AVAILABLE = True
except ImportError:
    PYMUPDF_AVAILABLE = False
    fitz = None

try:
    import textract
    TEXTRACT_AVAILABLE = True
except ImportError:
    TEXTRACT_AVAILABLE = False
    textract = None

try:
    import boto3
    AWS_AVAILABLE = True
except ImportError:
    AWS_AVAILABLE = False
    boto3 = None

try:
    from google.cloud import bigquery
    BIGQUERY_AVAILABLE = True
except ImportError:
    BIGQUERY_AVAILABLE = False
    bigquery = None

try:
    import psycopg2
    PSYCOPG2_AVAILABLE = True
except ImportError:
    PSYCOPG2_AVAILABLE = False
    psycopg2 = None

try:
    import mysql.connector
    MYSQL_AVAILABLE = True
except ImportError:
    MYSQL_AVAILABLE = False
    mysql = None

try:
    import pymongo
    MONGODB_AVAILABLE = True
except ImportError:
    MONGODB_AVAILABLE = False
    pymongo = None

try:
    from sqlalchemy import create_engine, text
    SQLALCHEMY_AVAILABLE = True
except ImportError:
    SQLALCHEMY_AVAILABLE = False
    create_engine = None
    text = None

try:
    import openpyxl
    OPENPYXL_AVAILABLE = True
except ImportError:
    OPENPYXL_AVAILABLE = False
    openpyxl = None

try:
    from PIL import Image
    PIL_AVAILABLE = True
except ImportError:
    PIL_AVAILABLE = False
    Image = None

try:
    import pytesseract
    TESSERACT_AVAILABLE = True
except ImportError:
    TESSERACT_AVAILABLE = False
    pytesseract = None

try:
    import asyncpg
    ASYNCPG_AVAILABLE = True
except ImportError:
    ASYNCPG_AVAILABLE = False
    asyncpg = None

try:
    import duckdb
    DUCKDB_AVAILABLE = True
except ImportError:
    DUCKDB_AVAILABLE = False
    duckdb = None

# Load environment variables
load_dotenv()

warnings.filterwarnings('ignore')
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

app = Flask(__name__)

# OCR API Configuration
OCR_SPACE_API_KEY = os.getenv('OCR_SPACE_API_KEY', 'K87899142688957')  # Free API key
OCR_SPACE_URL = 'https://api.ocr.space/parse/image'

# Custom JSON encoder to handle NumPy/Pandas types
class NumpyJSONEncoder(json.JSONEncoder):
    def default(self, obj):
        if isinstance(obj, (np.integer, np.int64, np.int32)):
            return int(obj)
        elif isinstance(obj, (np.floating, np.float64, np.float32)):
            return float(obj)
        elif isinstance(obj, np.ndarray):
            return obj.tolist()
        elif isinstance(obj, pd.Timestamp):
            return obj.isoformat()
        elif isinstance(obj, pd.Series):
            return obj.tolist()
        elif pd.isna(obj):
            return None
        return super().default(obj)

# Timeout decorator for function-level timeouts - simplified version
@contextmanager
def timeout_context(seconds):
    """Simplified timeout context that doesn't use signals"""
    # For now, just yield without timeout to avoid signal issues
    # TODO: Implement proper timeout handling for production
    yield

class UniversalDataAnalystAgent:
    def __init__(self):
        # Initialize all possible API clients
        self.anthropic = Anthropic(api_key=os.getenv('ANTHROPIC_API_KEY'))
        self.max_retries = 2
        self.start_time = None
        self.timeout_seconds = 170
        
        # Database connections pool
        self.db_connections = {}
        
        # Enhanced file processing extensions
        self.supported_extensions = {
            'tabular': ['.csv', '.tsv', '.xlsx', '.xls', '.parquet', '.json', '.jsonl', '.xml', '.html'],
            'text': ['.txt', '.md', '.rtf', '.doc', '.docx'],
            'image': ['.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff', '.webp', '.svg'],
            'pdf': ['.pdf'],
            'archive': ['.zip', '.tar', '.gz', '.7z', '.rar'],
            'database': ['.db', '.sqlite', '.sqlite3'],
            'code': ['.py', '.js', '.sql', '.r', '.scala', '.java', '.cpp', '.c', '.h'],
            'presentation': ['.ppt', '.pptx'],
            'other': ['.log', '.conf', '.ini', '.yaml', '.yml']
        }
        
        # Initialize OCR and vision capabilities
        self.init_ocr_capabilities()
        
        # Initialize cloud clients if credentials available
        self.init_cloud_clients()
    
    def init_ocr_capabilities(self):
        """Initialize OCR capabilities - handle missing dependencies gracefully"""
        try:
            if EASYOCR_AVAILABLE:
                self.easyocr_reader = easyocr.Reader(['en'])
            else:
                self.easyocr_reader = None
                logger.warning("EasyOCR not installed - OCR fallback disabled")
        except Exception as e:
            logger.warning(f"EasyOCR initialization failed: {e}")
            self.easyocr_reader = None
    
    def _is_valid_dataframe(self, df) -> bool:
        """Check if DataFrame is valid and has data"""
        return df is not None and isinstance(df, pd.DataFrame) and len(df) > 0
    
    def _is_empty_or_none(self, df) -> bool:
        """Check if DataFrame is None or empty"""
        return df is None or not isinstance(df, pd.DataFrame) or len(df) == 0
        
    def init_cloud_clients(self):
        """Initialize cloud service clients if credentials are available"""
        try:
            # AWS
            if os.getenv('AWS_ACCESS_KEY_ID') and AWS_AVAILABLE:
                self.s3_client = boto3.client('s3')
                self.textract_client = boto3.client('textract')
            
            # Google Cloud
            if os.getenv('GOOGLE_APPLICATION_CREDENTIALS') and BIGQUERY_AVAILABLE:
                self.bigquery_client = bigquery.Client()
        except Exception as e:
            logger.warning(f"Could not initialize cloud clients: {e}")
    
    def call_claude_advanced(self, prompt: str, max_tokens: int = 4000, include_images: List[bytes] = None) -> str:
        """Enhanced Claude API call with FIXED image media type detection"""
        try:
            with timeout_context(30):
                messages = []
                
                if include_images:
                    content = [{"type": "text", "text": prompt}]
                    for img_bytes in include_images[:3]:
                        img_b64 = base64.b64encode(img_bytes).decode('utf-8')
                        media_type = self._detect_image_media_type(img_bytes)
                        
                        content.append({
                            "type": "image",
                            "source": {
                                "type": "base64",
                                "media_type": media_type,
                                "data": img_b64
                            }
                        })
                    messages.append({"role": "user", "content": content})
                else:
                    messages.append({"role": "user", "content": prompt})
                
                response = self.anthropic.messages.create(
                    model="claude-3-5-sonnet-20241022",
                    max_tokens=max_tokens,
                    temperature=0,
                    messages=messages
                )
                return response.content[0].text
        except Exception as e:
            logger.error(f"Claude API error: {e}")
            raise e

    def _detect_image_media_type(self, img_bytes: bytes) -> str:
        """Detect actual image media type from bytes"""
        try:
            if img_bytes.startswith(b'\x89PNG\r\n\x1a\n'):
                return "image/png"
            elif img_bytes.startswith(b'\xff\xd8\xff'):
                return "image/jpeg"
            elif img_bytes.startswith(b'GIF8'):
                return "image/gif"
            elif img_bytes.startswith(b'RIFF') and b'WEBP' in img_bytes[:12]:
                return "image/webp"
            elif img_bytes.startswith(b'BM'):
                return "image/bmp"
            elif img_bytes.startswith(b'II*\x00') or img_bytes.startswith(b'MM\x00*'):
                return "image/tiff"
            else:
                if PIL_AVAILABLE:
                    try:
                        image = Image.open(BytesIO(img_bytes))
                        format_lower = image.format.lower()
                        return f"image/{format_lower}"
                    except:
                        pass
                return "image/jpeg"
        except:
            return "image/jpeg"

    def time_remaining(self) -> float:
        """Check remaining time"""
        if not self.start_time:
            return self.timeout_seconds
        elapsed = time.time() - self.start_time
        return max(0, self.timeout_seconds - elapsed)

    def detect_data_sources(self, question: str, files: Dict = None) -> Dict[str, Any]:
        """Universal data source detection - optimized for speed"""
        # Clean URLs - fix malformed URLs
        url_pattern = r'https?://[^\s\'"<>)]+(?:\)[^\s\'"<>]*)?'
        raw_urls = re.findall(url_pattern, question)
        cleaned_urls = []
        
        for url in raw_urls:
            # Remove trailing punctuation that's not part of URL
            cleaned_url = re.sub(r'[.)]+$', '', url)
            if cleaned_url.endswith('/)'):
                cleaned_url = cleaned_url[:-1]  # Remove trailing )
            cleaned_urls.append(cleaned_url)
        
        sources = {
            'urls': cleaned_urls,
            'ftp_urls': re.findall(r'ftp://[^\s\'"<>]+', question),
            's3_paths': re.findall(r's3://[^\s\'"<>]+', question),
            'gcs_paths': re.findall(r'gs://[^\s\'"<>]+', question),
            'database_urls': re.findall(r'(postgresql|mysql|mongodb|sqlite|oracle|mssql)://[^\s\'"<>]+', question),
            'api_endpoints': re.findall(r'(api|endpoint|service)\s*:?\s*https?://[^\s\'"<>]+', question, re.IGNORECASE),
            'file_patterns': re.findall(r'\*\.[a-zA-Z0-9]+', question),
            'table_names': re.findall(r'(table|collection|dataset)\s+([a-zA-Z_][a-zA-Z0-9_]*)', question, re.IGNORECASE),
            'query_hints': re.findall(r'(SELECT|FROM|WHERE|JOIN|GROUP BY|ORDER BY)', question, re.IGNORECASE),
            'attached_files': list(files.keys()) if files else [],
            'cloud_references': {
                'aws': bool(re.search(r'(aws|amazon|s3|ec2|rds|redshift)', question, re.IGNORECASE)),
                'gcp': bool(re.search(r'(google cloud|gcp|bigquery|cloud storage)', question, re.IGNORECASE)),
                'azure': bool(re.search(r'(azure|microsoft cloud)', question, re.IGNORECASE))
            }
        }
        
        # Extract context clues
        sources['context'] = {
            'needs_scraping': bool(sources['urls'] or re.search(r'(scrape|extract from|download from)', question, re.IGNORECASE)),
            'needs_plotting': bool(re.search(r'(plot|graph|chart|visualiz|scatter|bar|line|histogram)', question, re.IGNORECASE)),
            'needs_analysis': bool(re.search(r'(analyz|correlat|regression|predict|trend|pattern)', question, re.IGNORECASE)),
            'needs_aggregation': bool(re.search(r'(count|sum|average|mean|median|group|aggregate)', question, re.IGNORECASE)),
            'time_series': bool(re.search(r'(time|date|year|month|day|temporal|trend over)', question, re.IGNORECASE)),
            'json_output': bool(re.search(r'(json|JSON)', question, re.IGNORECASE)),
            'array_output': bool(re.search(r'(array|list of)', question, re.IGNORECASE))
        }
        
        logger.info(f"Detected sources: {sources}")
        return sources

    async def universal_web_scraper(self, url: str, max_attempts: int = 2) -> str:
        """Universal web scraper that handles any website - optimized for speed"""
        logger.info(f"Scraping: {url}")
        
        # Try methods in order of preference - reduced for speed
        methods = [
            self._scrape_with_requests,
            self._scrape_with_aiohttp if AIOHTTP_AVAILABLE else None,
            self._scrape_with_playwright if PLAYWRIGHT_AVAILABLE else None
        ]
        
        # Filter out None methods
        methods = [method for method in methods if method is not None]
        
        for i, method in enumerate(methods):
            try:
                if self.time_remaining() < 20:
                    break
                    
                logger.info(f"Trying scraping method {i+1}/{len(methods)}")
                with timeout_context(25):
                    content = await method(url)
                    if content and len(content) > 100:
                        logger.info(f"Successful scrape with method {i+1}, content length: {len(content)}")
                        return content
            except Exception as e:
                logger.warning(f"Scraping method {i+1} failed: {e}")
                continue
        
        raise Exception(f"All scraping methods failed for {url}")

    async def _scrape_with_requests(self, url: str) -> str:
        """Standard requests scraping"""
        headers = {
            'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36',
            'Accept': 'text/html,application/xhtml+xml,application/xml;q=0.9,image/webp,*/*;q=0.8',
            'Accept-Language': 'en-US,en;q=0.5',
            'Accept-Encoding': 'gzip, deflate',
            'Connection': 'keep-alive',
            'Upgrade-Insecure-Requests': '1'
        }
        
        session = requests.Session()
        session.headers.update(headers)
        
        response = session.get(url, timeout=20, allow_redirects=True, verify=False)
        response.raise_for_status()
        
        # Handle encoding
        if response.encoding is None:
            detected = chardet.detect(response.content[:1000])
            response.encoding = detected.get('encoding', 'utf-8')
        
        return response.text

    async def _scrape_with_aiohttp(self, url: str) -> str:
        """Async HTTP scraping"""
        if not AIOHTTP_AVAILABLE:
            raise Exception("aiohttp not available")
            
        headers = {
            'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36'
        }
        
        timeout = aiohttp.ClientTimeout(total=20)
        async with aiohttp.ClientSession(timeout=timeout, headers=headers) as session:
            async with session.get(url, ssl=False) as response:
                content = await response.read()
                detected = chardet.detect(content[:1000])
                encoding = detected.get('encoding', 'utf-8')
                return content.decode(encoding, errors='ignore')

    async def _scrape_with_playwright(self, url: str) -> str:
        """Dynamic content scraping with Playwright"""
        if not PLAYWRIGHT_AVAILABLE:
            raise Exception("Playwright not available")
            
        async with async_playwright() as p:
            browser = await p.chromium.launch(
                headless=True,
                args=['--no-sandbox', '--disable-dev-shm-usage', '--disable-gpu', '--disable-images']
            )
            
            context = await browser.new_context(
                user_agent='Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36',
                viewport={'width': 1280, 'height': 720}
            )
            
            page = await context.new_page()
            
            try:
                await page.goto(url, wait_until='domcontentloaded', timeout=30000)
                await page.wait_for_timeout(1000)
                content = await page.content()
                return content
                
            finally:
                await browser.close()

    def universal_data_extractor(self, content: str, source_url: str = None) -> pd.DataFrame:
        """Extract structured data from any content type - optimized"""
        logger.info("Starting universal data extraction")
        
        # Determine content type
        content_type = self._detect_content_type(content, source_url)
        
        extractors = {
            'html': self._extract_from_html,
            'json': self._extract_from_json,
            'xml': self._extract_from_xml,
            'csv': self._extract_from_csv,
            'text': self._extract_from_text,
            'api': self._extract_from_api_response
        }
        
        extractor = extractors.get(content_type, self._extract_from_html)
        
        try:
            with timeout_context(30):
                df = extractor(content)
                if self._is_valid_dataframe(df):
                    logger.info(f"Extracted {df.shape[0]} rows, {df.shape[1]} columns")
                    return self._clean_dataframe(df)
        except Exception as e:
            logger.warning(f"Primary extraction failed: {e}")
        
        # Fallback: try most common extractors only
        for name, extractor in [('json', extractors['json']), ('csv', extractors['csv'])]:
            try:
                with timeout_context(15):
                    df = extractor(content)
                    if self._is_valid_dataframe(df):
                        logger.info(f"Fallback extraction successful with {name}")
                        return self._clean_dataframe(df)
            except:
                continue
        
        # LLM-based extraction if time permits
        if self.time_remaining() > 30:
            return self._extract_with_llm_analysis(content)
        
        return pd.DataFrame()

    def _detect_content_type(self, content: str, source_url: str = None) -> str:
        """Detect the type of content"""
        content_lower = content.lower().strip()
        
        if source_url:
            if any(api_term in source_url.lower() for api_term in ['api', 'json', 'rest']):
                return 'api'
        
        if content_lower.startswith('<!doctype html') or content_lower.startswith('<html'):
            return 'html'
        elif content_lower.startswith('{') or content_lower.startswith('['):
            return 'json'
        elif content_lower.startswith('<?xml') or '<xml' in content_lower[:100]:
            return 'xml'
        elif ',' in content and '\n' in content and content.count(',') > content.count(' '):
            return 'csv'
        else:
            return 'text'

    def _extract_from_html(self, html_content: str) -> pd.DataFrame:
        """Enhanced HTML extraction - optimized"""
        soup = BeautifulSoup(html_content, 'html.parser')
        
        # Strategy 1: Tables (most reliable) - limit to first 5 tables
        tables = soup.find_all('table')[:5]
        for table in tables:
            try:
                df = self._extract_table_data(table)
                if self._is_valid_dataframe(df) and len(df) > 2:
                    return df
            except:
                continue
        
        # Strategy 2: Structured lists - limit to first 3 lists
        lists = soup.find_all(['ul', 'ol'])[:3]
        for list_elem in lists:
            try:
                df = self._extract_list_data(list_elem)
                if self._is_valid_dataframe(df) and len(df) > 2:
                    return df
            except:
                continue
        
        # Strategy 3: JSON-LD structured data
        json_scripts = soup.find_all('script', type='application/ld+json')[:2]
        for script in json_scripts:
            try:
                data = json.loads(script.string)
                df = pd.json_normalize(data)
                if len(df) > 0:
                    return df
            except:
                continue
        
        return None

    def _extract_table_data(self, table) -> pd.DataFrame:
        """Extract data from HTML table"""
        rows = table.find_all('tr')
        if len(rows) < 2:
            return None
        
        # Get headers
        header_row = rows[0]
        headers = []
        for cell in header_row.find_all(['th', 'td']):
            text = self._clean_text(cell.get_text())
            headers.append(text if text else f'Column_{len(headers)}')
        
        if not headers:
            return None
        
        # Get data rows - limit to first 1000 rows for speed
        data_rows = []
        for row in rows[1:1001]:
            cells = row.find_all(['td', 'th'])
            if not cells:
                continue
            
            row_data = []
            for cell in cells:
                text = self._clean_text(cell.get_text())
                row_data.append(text)
            
            # Ensure row matches header length
            while len(row_data) < len(headers):
                row_data.append('')
            row_data = row_data[:len(headers)]
            
            if any(cell.strip() for cell in row_data):
                data_rows.append(row_data)
        
        if not data_rows:
            return None
        
        return pd.DataFrame(data_rows, columns=headers)

    def _extract_list_data(self, list_elem) -> pd.DataFrame:
        """Extract structured data from lists"""
        items = list_elem.find_all('li')
        if len(items) < 3:
            return None
        
        data_rows = []
        for item in items[:100]:  # Limit for speed
            text = self._clean_text(item.get_text())
            
            # Try different separators
            for sep in [':', '|', '\t', ' - ', ' — ']:
                if sep in text:
                    parts = [part.strip() for part in text.split(sep)]
                    if len(parts) >= 2 and all(part for part in parts):
                        data_rows.append(parts)
                        break
        
        if len(data_rows) < 3:
            return None
        
        # Normalize row lengths
        max_cols = max(len(row) for row in data_rows)
        headers = [f'Field_{i}' for i in range(max_cols)]
        
        for row in data_rows:
            while len(row) < max_cols:
                row.append('')
        
        return pd.DataFrame(data_rows, columns=headers)

    def _extract_from_json(self, content: str) -> pd.DataFrame:
        """Extract data from JSON content"""
        try:
            data = json.loads(content)
            
            if isinstance(data, list):
                return pd.json_normalize(data)
            elif isinstance(data, dict):
                # Look for array properties
                for key, value in data.items():
                    if isinstance(value, list) and len(value) > 2:
                        return pd.json_normalize(value)
                
                # If no arrays, normalize the dict itself
                return pd.json_normalize([data])
            
        except json.JSONDecodeError:
            # Try to extract JSON from text
            json_patterns = re.findall(r'\{[^{}]*\}|\[[^\[\]]*\]', content)[:5]
            for pattern in json_patterns:
                try:
                    data = json.loads(pattern)
                    df = pd.json_normalize(data)
                    if len(df) > 0:
                        return df
                except:
                    continue
        
        return None

    def _extract_from_xml(self, content: str) -> pd.DataFrame:
        """Extract data from XML content"""
        try:
            import xml.etree.ElementTree as ET
            root = ET.fromstring(content)
            
            # Find repeating elements
            tag_counts = {}
            for elem in root.iter():
                tag_counts[elem.tag] = tag_counts.get(elem.tag, 0) + 1
            
            # Get the most common tag (likely data rows)
            most_common_tag = max(tag_counts, key=tag_counts.get)
            
            if tag_counts[most_common_tag] > 2:
                elements = root.findall(f'.//{most_common_tag}')[:200]
                
                data_rows = []
                for elem in elements:
                    row_data = {}
                    for child in elem:
                        row_data[child.tag] = child.text or ''
                    
                    # Also include attributes
                    row_data.update(elem.attrib)
                    
                    if row_data:
                        data_rows.append(row_data)
                
                if data_rows:
                    return pd.DataFrame(data_rows)
        
        except Exception as e:
            logger.warning(f"XML parsing failed: {e}")
        
        return None

    def _extract_from_csv(self, content: str) -> pd.DataFrame:
        """Extract data from CSV-like content"""
        # Try different separators
        separators = [',', '\t', '|', ';']
        
        for sep in separators:
            try:
                df = pd.read_csv(StringIO(content), sep=sep, error_bad_lines=False, nrows=10000)
                if len(df) > 2 and len(df.columns) > 1:
                    return df
            except:
                continue
        
        return None

    def _extract_from_text(self, content: str) -> pd.DataFrame:
        """Extract structured data from plain text"""
        lines = content.split('\n')[:200]
        
        # Look for patterns that suggest tabular data
        structured_lines = []
        for line in lines:
            line = line.strip()
            if not line:
                continue
            
            # Check for separators
            for sep in ['\t', '|', ':', ' - ', '  ']:
                if sep in line and len(line.split(sep)) > 1:
                    parts = [part.strip() for part in line.split(sep)]
                    if len(parts) >= 2 and all(part for part in parts):
                        structured_lines.append(parts)
                        break
        
        if len(structured_lines) > 2:
            max_cols = max(len(row) for row in structured_lines)
            headers = [f'Column_{i}' for i in range(max_cols)]
            
            for row in structured_lines:
                while len(row) < max_cols:
                    row.append('')
            
            return pd.DataFrame(structured_lines, columns=headers)
        
        return None

    def _extract_from_api_response(self, content: str) -> pd.DataFrame:
        """Extract data from API response"""
        # Try JSON first
        df = self._extract_from_json(content)
        if df is not None:
            return df
        
        # Try XML
        df = self._extract_from_xml(content)
        if df is not None:
            return df
        
        return None

    def _extract_with_llm_analysis(self, content: str) -> pd.DataFrame:
        """Use LLM to analyze and extract data from any content"""
        if self.time_remaining() < 25:
            return pd.DataFrame()
        
        analysis_prompt = f"""
Analyze this content and extract any structured data you can find. Return as JSON array of objects.

Content (first 2000 chars):
{content[:2000]}

Look for:
- Tables with headers and data rows
- Lists of similar items
- Key-value pairs that repeat
- Any patterns that could be structured data

Return ONLY a JSON array like: [{{"field1": "value1", "field2": "value2"}}, {{"field1": "value3", "field2": "value4"}}]

If no structured data found, return empty array: []
"""
        
        try:
            response = self.call_claude_advanced(analysis_prompt, 2000)
            
            # Extract JSON from response
            json_match = re.search(r'\[.*\]', response, re.DOTALL)
            if json_match:
                data = json.loads(json_match.group())
                if isinstance(data, list) and len(data) > 0:
                    return pd.DataFrame(data)
        
        except Exception as e:
            logger.warning(f"LLM-based extraction failed: {e}")
        
        return pd.DataFrame()

    def _clean_dataframe(self, df: pd.DataFrame) -> pd.DataFrame:
        """Clean and standardize DataFrame - with better error handling"""
        if df is None or len(df) == 0:
            logger.warning("_clean_dataframe received None or empty DataFrame")
            return pd.DataFrame()
        
        logger.info(f"Cleaning DataFrame: {df.shape}")
        
        try:
            # Create a copy to avoid modifying original
            df_clean = df.copy()
            
            # Clean column names
            df_clean.columns = [self._clean_text(str(col)) for col in df_clean.columns]
            
            # Remove completely empty rows and columns
            original_shape = df_clean.shape
            df_clean = df_clean.dropna(how='all').dropna(axis=1, how='all')
            logger.info(f"After removing empty rows/cols: {original_shape} -> {df_clean.shape}")
            
            if len(df_clean) == 0:
                logger.warning("DataFrame became empty after cleaning")
                return pd.DataFrame()
            
            # Clean text in object columns - limit to first 1000 rows for speed
            sample_df = df_clean.head(1000) if len(df_clean) > 1000 else df_clean
            for col in sample_df.select_dtypes(include=['object']).columns:
                try:
                    df_clean[col] = df_clean[col].astype(str).apply(lambda x: self._clean_text(x) if pd.notna(x) else x)
                except Exception as e:
                    logger.warning(f"Failed to clean column {col}: {e}")
                    continue
            
            # Try to convert columns to appropriate types with better error handling
            for col in df_clean.columns:
                if df_clean[col].dtype == 'object':
                    try:
                        # Try numeric conversion on sample
                        sample_series = df_clean[col].head(100)
                        # Remove non-numeric strings before conversion
                        clean_sample = sample_series.astype(str).str.replace(r'[^\d.-]', '', regex=True)
                        numeric_series = pd.to_numeric(clean_sample, errors='coerce')
                        
                        if not numeric_series.isna().all() and numeric_series.notna().sum() / len(sample_series) > 0.5:
                            # Apply to full column
                            clean_col = df_clean[col].astype(str).str.replace(r'[^\d.-]', '', regex=True)
                            df_clean[col] = pd.to_numeric(clean_col, errors='coerce')
                    except Exception as e:
                        logger.debug(f"Failed to convert column {col} to numeric: {e}")
                        continue
            
            logger.info(f"Final cleaned DataFrame: {df_clean.shape}")
            return df_clean
            
        except Exception as e:
            logger.error(f"DataFrame cleaning failed: {e}")
            return df  # Return original if cleaning fails

    def _clean_text(self, text: str) -> str:
        """Clean and normalize text with better error handling"""
        if not isinstance(text, str):
            text = str(text)
        
        try:
            # Remove control characters
            text = ''.join(char for char in text if unicodedata.category(char)[0] != 'C')
            
            # Normalize whitespace
            text = re.sub(r'\s+', ' ', text.strip())
            
            # Remove common artifacts
            text = re.sub(r'\[.*?\]', '', text)
            text = re.sub(r'\(edit\)', '', text, flags=re.IGNORECASE)
            
            return text.strip()
        except Exception as e:
            logger.debug(f"Text cleaning failed: {e}")
            return str(text) if text is not None else ""

    def universal_database_handler(self, question: str, sources: Dict) -> Any:
        """Handle any type of database or data source - optimized"""
        logger.info("Processing database/file sources")
        
        if self.time_remaining() < 30:
            return {"error": "Insufficient time for database operations"}
        
        try:
            if sources.get('s3_paths'):
                with timeout_context(45):
                    return self._handle_s3_sources(sources['s3_paths'], question)
            
            if sources.get('database_urls'):
                with timeout_context(30):
                    return self._handle_database_urls(sources['database_urls'], question)
            
            if sources.get('query_hints'):
                with timeout_context(20):
                    return self._handle_sql_queries(question, sources)
        
        except TimeoutError as e:
            logger.error(f"Database operation timed out: {e}")
            return {"error": "Database operation timed out"}
        except Exception as e:
            logger.error(f"Database operation failed: {e}")
            return {"error": f"Database operation failed: {str(e)}"}
        
        return {"error": "No database sources found"}

    def _handle_s3_sources(self, s3_paths: List[str], question: str) -> Any:
        """Handle S3 data sources using LLM to generate queries"""
        try:
            if not DUCKDB_AVAILABLE:
                return {"error": "DuckDB not available for S3 processing"}
                
            conn = duckdb.connect()
            
            # Install and load extensions
            extensions = ['httpfs', 'parquet', 'json', 'csv']
            for ext in extensions:
                try:
                    conn.execute(f"INSTALL {ext}; LOAD {ext};")
                except:
                    pass
            
            # Set S3 configuration
            if any('ap-south-1' in path for path in s3_paths):
                conn.execute("SET s3_region='ap-south-1';")
            
            # Use LLM to generate appropriate SQL query
            query_prompt = f"""
Generate a DuckDB SQL query to answer this question: {question}

Available S3 paths: {s3_paths}

Guidelines:
- Use read_parquet() for .parquet files
- Use read_csv() for .csv files
- Use appropriate aggregations and filters based on the question
- Include proper GROUP BY, ORDER BY as needed
- Handle date/time columns appropriately

Return ONLY the SQL query, no explanations:
"""
            
            sql_query = self.call_claude_advanced(query_prompt, 1500)
            sql_query = re.sub(r'```sql|```', '', sql_query).strip()
            
            logger.info(f"Generated SQL: {sql_query}")
            
            result = conn.execute(sql_query).fetchall()
            conn.close()
            
            return result
            
        except Exception as e:
            logger.error(f"S3 query failed: {e}")
            return {"error": f"S3 query failed: {str(e)}"}

    def _handle_database_urls(self, db_urls: List[str], question: str) -> Any:
        """Handle direct database connections using LLM-generated queries"""
        for url in db_urls[:2]:
            try:
                with timeout_context(25):
                    if url.startswith('postgresql://'):
                        return self._query_postgresql(url, question)
                    elif url.startswith('mysql://'):
                        return self._query_mysql(url, question)
                    elif url.startswith('sqlite://'):
                        return self._query_sqlite(url, question)
            except Exception as e:
                logger.error(f"Database query failed for {url}: {e}")
                continue
        
        return {"error": "All database connections failed"}

    def _query_postgresql(self, url: str, question: str) -> Any:
        """Query PostgreSQL database with LLM-generated SQL"""
        try:
            if not SQLALCHEMY_AVAILABLE:
                return {"error": "SQLAlchemy not available for PostgreSQL"}
                
            engine = create_engine(url, connect_args={"connect_timeout": 15})
            
            # Use LLM to generate SQL query
            sql_query = self._generate_sql_from_question(question, 'postgresql')
            
            with engine.connect() as conn:
                result = conn.execute(text(sql_query)).fetchall()
                return [dict(row._mapping) for row in result]
        
        except Exception as e:
            logger.error(f"PostgreSQL query failed: {e}")
            return {"error": f"PostgreSQL query failed: {str(e)}"}

    def _query_mysql(self, url: str, question: str) -> Any:
        """Query MySQL database with LLM-generated SQL"""
        try:
            if not SQLALCHEMY_AVAILABLE:
                return {"error": "SQLAlchemy not available for MySQL"}
                
            engine = create_engine(url, connect_args={"connect_timeout": 15})
            sql_query = self._generate_sql_from_question(question, 'mysql')
            
            with engine.connect() as conn:
                result = conn.execute(text(sql_query)).fetchall()
                return [dict(row._mapping) for row in result]
        
        except Exception as e:
            logger.error(f"MySQL query failed: {e}")
            return {"error": f"MySQL query failed: {str(e)}"}

    def _query_sqlite(self, url: str, question: str) -> Any:
        """Query SQLite database with LLM-generated SQL"""
        try:
            db_path = url.replace('sqlite://', '')
            conn = sqlite3.connect(db_path, timeout=15)
            
            sql_query = self._generate_sql_from_question(question, 'sqlite')
            
            cursor = conn.execute(sql_query)
            result = cursor.fetchall()
            columns = [description[0] for description in cursor.description]
            
            conn.close()
            
            return [dict(zip(columns, row)) for row in result]
        
        except Exception as e:
            logger.error(f"SQLite query failed: {e}")
            return {"error": f"SQLite query failed: {str(e)}"}

    def _generate_sql_from_question(self, question: str, db_type: str) -> str:
        """Generate SQL query from natural language question using LLM"""
        prompt = f"""
Generate a {db_type.upper()} SQL query to answer this question: {question}

Requirements:
- Use appropriate table names and column names based on the question context
- Include proper JOINs, WHERE clauses, GROUP BY, ORDER BY as needed
- Handle date/time operations appropriately
- Use appropriate aggregate functions (COUNT, SUM, AVG, etc.)
- Limit results if needed for performance

Return ONLY the SQL query:
"""
        
        response = self.call_claude_advanced(prompt, 1000)
        return re.sub(r'```sql|```', '', response).strip()

    def _handle_sql_queries(self, question: str, sources: Dict) -> Any:
        """Handle raw SQL queries in the question using LLM enhancement"""
        sql_patterns = re.findall(r'```sql\s*(.*?)\s*```', question, re.DOTALL | re.IGNORECASE)
        if not sql_patterns:
            sql_patterns = re.findall(r'(SELECT.*?;?)\s*', question, re.DOTALL | re.IGNORECASE | re.MULTILINE)
        
        if sql_patterns:
            try:
                if not DUCKDB_AVAILABLE:
                    return {"error": "DuckDB not available for SQL processing"}
                    
                conn = duckdb.connect()
                
                # Install extensions
                for ext in ['httpfs', 'parquet', 'json', 'csv']:
                    try:
                        conn.execute(f"INSTALL {ext}; LOAD {ext};")
                    except:
                        pass
                
                sql_query = sql_patterns[0].strip()
                
                # Enhance query with LLM if needed
                if "?" in question or "enhance" in question.lower():
                    enhancement_prompt = f"""
Enhance this SQL query to better answer the user's question: {question}

Original query: {sql_query}

Improvements to consider:
- Add missing columns or calculations
- Improve filtering or grouping
- Add proper ordering
- Handle edge cases

Return the enhanced SQL query:
"""
                    enhanced_query = self.call_claude_advanced(enhancement_prompt, 1000)
                    sql_query = re.sub(r'```sql|```', '', enhanced_query).strip()
                
                result = conn.execute(sql_query).fetchall()
                conn.close()
                
                return result
                
            except Exception as e:
                logger.error(f"SQL execution failed: {e}")
                return {"error": f"SQL execution failed: {str(e)}"}
        
        return {"error": "No valid SQL found"}

    def process_attached_files(self, files: Dict) -> Dict[str, Any]:
        """Process any type of attached file - comprehensive support"""
        processed = {}
        
        for filename, file_obj in files.items():
            if self.time_remaining() < 15:
                break
                
            try:
                logger.info(f"Processing file: {filename}")
                
                # Reset file pointer
                file_obj.seek(0)
                
                # Get file extension
                _, ext = os.path.splitext(filename.lower())
                logger.info(f"File extension: {ext}")
                
                # Process based on file type
                try:
                    if ext in ['.csv', '.tsv']:
                        logger.info(f"Processing as CSV/TSV: {filename}")
                        df = self._process_csv_file(file_obj, ext)
                        if self._is_valid_dataframe(df):
                            logger.info(f"Successfully processed CSV: {df.shape}")
                            processed[filename] = df
                        else:
                            logger.warning(f"CSV processing returned invalid DataFrame")
                            processed[filename] = {"error": "Invalid DataFrame from CSV processing"}
                    
                    elif ext in ['.xlsx', '.xls']:
                        logger.info(f"Processing as Excel: {filename}")
                        df = self._process_excel_file(file_obj)
                        processed[filename] = df
                    
                    elif ext == '.json':
                        logger.info(f"Processing as JSON: {filename}")
                        data = self._process_json_file(file_obj)
                        processed[filename] = data
                    
                    elif ext == '.parquet':
                        logger.info(f"Processing as Parquet: {filename}")
                        df = self._process_parquet_file(file_obj)
                        processed[filename] = df
                    
                    elif ext in ['.txt', '.md', '.log']:
                        logger.info(f"Processing as text: {filename}")
                        text = self._process_text_file(file_obj)
                        processed[filename] = text
                    
                    elif ext == '.pdf':
                        logger.info(f"Processing as PDF: {filename}")
                        text = self._process_pdf_file(file_obj)
                        processed[filename] = text
                    
                    elif ext in ['.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff', '.webp']:
                        logger.info(f"Processing as image: {filename}")
                        result = self._process_image_file(file_obj, filename)
                        processed[filename] = result
                    
                    elif ext in ['.doc', '.docx']:
                        logger.info(f"Processing as Word document: {filename}")
                        text = self._process_word_file(file_obj)
                        processed[filename] = text
                    
                    elif ext in ['.ppt', '.pptx']:
                        logger.info(f"Processing as PowerPoint: {filename}")
                        text = self._process_powerpoint_file(file_obj)
                        processed[filename] = text
                    
                    elif ext == '.zip':
                        logger.info(f"Processing as ZIP: {filename}")
                        extracted = self._process_zip_file(file_obj)
                        processed.update(extracted)
                    
                    elif ext in ['.db', '.sqlite', '.sqlite3']:
                        logger.info(f"Processing as SQLite database: {filename}")
                        tables = self._process_sqlite_file(file_obj)
                        processed[filename] = tables
                    
                    elif ext in ['.xml', '.html']:
                        logger.info(f"Processing as XML/HTML: {filename}")
                        data = self._process_xml_html_file(file_obj)
                        processed[filename] = data
                    
                    else:
                        # Try to read as CSV first
                        logger.info(f"Unknown extension {ext}, trying as CSV first: {filename}")
                        try:
                            file_obj.seek(0)
                            df = self._process_csv_file(file_obj, '.csv')
                            if self._is_valid_dataframe(df):
                                logger.info(f"Successfully processed unknown file as CSV: {df.shape}")
                                processed[filename] = df
                            else:
                                raise Exception("Not a valid CSV")
                        except Exception as csv_error:
                            # Try to read as text
                            logger.info(f"CSV attempt failed, trying as text: {filename}")
                            file_obj.seek(0)
                            content = file_obj.read()
                            if isinstance(content, bytes):
                                try:
                                    content = content.decode('utf-8')
                                except:
                                    content = content.decode('latin-1', errors='ignore')
                            processed[filename] = content
                
                except Exception as processing_error:
                    logger.error(f"Processing error for {filename}: {processing_error}")
                    processed[filename] = {"error": f"Processing failed: {str(processing_error)}"}
                
            except Exception as e:
                logger.error(f"Failed to process {filename}: {e}")
                processed[filename] = {"error": f"Failed to process: {str(e)}"}
        
        return processed

    def _process_csv_file(self, file_obj, ext: str) -> pd.DataFrame:
        """Process CSV/TSV files - enhanced with better error handling"""
        logger.info(f"Starting CSV processing for extension: {ext}")
        
        separator = '\t' if ext == '.tsv' else ','
        
        # Try different encodings
        encodings = ['utf-8', 'latin-1', 'cp1252', 'iso-8859-1']
        
        for encoding in encodings:
            try:
                file_obj.seek(0)
                content = file_obj.read()
                if isinstance(content, bytes):
                    content = content.decode(encoding, errors='ignore')
                
                logger.info(f"Trying encoding: {encoding}")
                
                # Try different separators if CSV
                if ext == '.csv':
                    for sep in [',', ';', '\t', '|']:
                        try:
                            logger.info(f"Trying separator: '{sep}'")
                            df = pd.read_csv(StringIO(content), sep=sep, nrows=20000, on_bad_lines='skip')
                            logger.info(f"Raw DataFrame created: shape={df.shape}")
                            
                            if len(df.columns) > 1 or (len(df.columns) == 1 and len(df) > 1):
                                logger.info(f"DataFrame meets criteria, cleaning...")
                                cleaned_df = self._clean_dataframe(df)
                                if cleaned_df is not None and len(cleaned_df) > 0:
                                    logger.info(f"Successfully processed CSV: {cleaned_df.shape}")
                                    return cleaned_df
                                else:
                                    logger.warning("Cleaning resulted in empty DataFrame")
                            else:
                                logger.info(f"DataFrame doesn't meet criteria: {len(df.columns)} columns, {len(df)} rows")
                        except Exception as e:
                            logger.debug(f"Failed with separator '{sep}': {e}")
                            continue
                else:
                    # TSV processing
                    try:
                        df = pd.read_csv(StringIO(content), sep=separator, nrows=20000, on_bad_lines='skip')
                        logger.info(f"TSV DataFrame created: shape={df.shape}")
                        if len(df) > 0:
                            cleaned_df = self._clean_dataframe(df)
                            if cleaned_df is not None and len(cleaned_df) > 0:
                                logger.info(f"Successfully processed TSV: {cleaned_df.shape}")
                                return cleaned_df
                    except Exception as e:
                        logger.debug(f"TSV processing failed: {e}")
                    
            except Exception as e:
                logger.debug(f"Failed with encoding '{encoding}': {e}")
                continue
        
        # Last resort: try pandas auto-detection
        logger.info("Trying pandas auto-detection...")
        try:
            file_obj.seek(0)
            df = pd.read_csv(file_obj, nrows=20000, on_bad_lines='skip')
            logger.info(f"Auto-detection DataFrame: shape={df.shape}")
            if len(df) > 0:
                cleaned_df = self._clean_dataframe(df)
                if cleaned_df is not None and len(cleaned_df) > 0:
                    logger.info(f"Successfully parsed with pandas auto-detection: {cleaned_df.shape}")
                    return cleaned_df
        except Exception as e:
            logger.debug(f"Pandas auto-detection failed: {e}")
        
        # Manual parsing fallback
        logger.info("Trying manual parsing...")
        try:
            file_obj.seek(0)
            content = file_obj.read()
            if isinstance(content, bytes):
                content = content.decode('utf-8', errors='ignore')
            
            lines = content.strip().split('\n')
            logger.info(f"Manual parsing: found {len(lines)} lines")
            
            if len(lines) > 1:
                # Detect separator from first line
                first_line = lines[0]
                separator = ','
                for sep in [',', '\t', ';', '|']:
                    if sep in first_line and len(first_line.split(sep)) > 1:
                        separator = sep
                        logger.info(f"Detected separator: '{separator}'")
                        break
                
                # Parse manually
                data = []
                headers = [col.strip().strip('"\'') for col in lines[0].split(separator)]
                logger.info(f"Headers: {headers}")
                
                for line in lines[1:]:
                    row = [cell.strip().strip('"\'') for cell in line.split(separator)]
                    # Pad row to match headers length
                    while len(row) < len(headers):
                        row.append('')
                    data.append(row[:len(headers)])
                
                if data:
                    df = pd.DataFrame(data, columns=headers)
                    logger.info(f"Manual DataFrame created: {df.shape}")
                    cleaned_df = self._clean_dataframe(df)
                    if cleaned_df is not None and len(cleaned_df) > 0:
                        logger.info(f"Successfully parsed manually: {cleaned_df.shape}")
                        return cleaned_df
        except Exception as e:
            logger.debug(f"Manual parsing failed: {e}")
        
        logger.error("All CSV parsing methods failed")
        raise Exception("Could not parse CSV file with any method")

    def _process_excel_file(self, file_obj) -> pd.DataFrame:
        """Process Excel files - enhanced"""
        try:
            # Try openpyxl first for .xlsx files
            try:
                excel_file = pd.ExcelFile(file_obj, engine='openpyxl')
            except:
                # Fall back to xlrd for .xls files
                excel_file = pd.ExcelFile(file_obj, engine='xlrd')
            
            # Get the largest sheet (most likely to contain data)
            largest_sheet = None
            max_rows = 0
            
            for sheet_name in excel_file.sheet_names[:5]:
                try:
                    df = pd.read_excel(file_obj, sheet_name=sheet_name, nrows=20000)
                    if len(df) > max_rows:
                        max_rows = len(df)
                        largest_sheet = df
                except Exception as e:
                    logger.warning(f"Failed to read sheet {sheet_name}: {e}")
                    continue
            
            return self._clean_dataframe(largest_sheet) if largest_sheet is not None else pd.DataFrame()
            
        except Exception as e:
            logger.error(f"Excel processing failed: {e}")
            return pd.DataFrame()

    def _process_json_file(self, file_obj) -> Any:
        """Process JSON files - enhanced"""
        try:
            file_obj.seek(0)
            content = file_obj.read()
            if isinstance(content, bytes):
                content = content.decode('utf-8')
            
            data = json.loads(content)
            
            # If it's a list of objects, convert to DataFrame
            if isinstance(data, list) and len(data) > 0:
                if isinstance(data[0], dict):
                    return pd.json_normalize(data)
                else:
                    # List of primitives
                    return pd.DataFrame({'values': data})
            elif isinstance(data, dict):
                # Try to find array values that could be tabular
                for key, value in data.items():
                    if isinstance(value, list) and len(value) > 0:
                        if isinstance(value[0], dict):
                            return pd.json_normalize(value)
                # If no arrays found, normalize the dict
                return pd.json_normalize([data])
            
            return data
            
        except Exception as e:
            logger.error(f"JSON processing failed: {e}")
            return {}

    def _process_parquet_file(self, file_obj) -> pd.DataFrame:
        """Process Parquet files"""
        try:
            df = pd.read_parquet(file_obj)
            return self._clean_dataframe(df)
        except Exception as e:
            logger.error(f"Parquet processing failed: {e}")
            return pd.DataFrame()

    def _process_text_file(self, file_obj) -> str:
        """Process text files - enhanced"""
        try:
            file_obj.seek(0)
            content = file_obj.read()
            if isinstance(content, bytes):
                # Try different encodings
                for encoding in ['utf-8', 'latin-1', 'cp1252', 'iso-8859-1']:
                    try:
                        return content.decode(encoding)
                    except:
                        continue
                return content.decode('utf-8', errors='ignore')
            return content
        except Exception as e:
            logger.error(f"Text processing failed: {e}")
            return ""

    def _process_pdf_file(self, file_obj) -> str:
        """Process PDF files - enhanced"""
        try:
            if not PYMUPDF_AVAILABLE:
                logger.warning("PyMuPDF not available for PDF processing")
                return "[PDF processing requires PyMuPDF - install with: pip install PyMuPDF]"
            
            file_obj.seek(0)
            pdf_bytes = file_obj.read()
            
            doc = fitz.open(stream=pdf_bytes, filetype="pdf")
            text = ""
            
            for page_num in range(min(50, len(doc))):  # Process up to 50 pages
                page = doc.load_page(page_num)
                text += page.get_text() + "\n"
            
            doc.close()
            return text
            
        except Exception as e:
            logger.error(f"PDF processing failed: {e}")
            return ""

    def _process_image_file(self, file_obj, filename: str) -> Dict[str, Any]:
        """Process image files using Claude Vision API as PRIMARY method, with fallbacks"""
        logger.info(f"Processing image file with Claude Vision: {filename}")
        
        try:
            file_obj.seek(0)
            img_bytes = file_obj.read()
            
            # PRIMARY METHOD: Use Claude Vision for comprehensive image analysis
            vision_prompt = """
Analyze this image comprehensively and provide detailed insights. Include:

1. **Visual Description**: What do you see in the image? Objects, people, scenes, text, etc.

2. **Text Content**: Extract any text, numbers, labels, or written content visible in the image

3. **Data & Structure**: If there are charts, graphs, tables, or structured data, describe and extract them

4. **Key Insights**: What are the main takeaways, patterns, or important information?

5. **Context & Meaning**: What is the purpose or context of this image?

6. **Inferences**: What can you infer or conclude from this image?

Provide a detailed analysis that would be useful for answering questions about this image.
"""
            
            logger.info("Attempting Claude Vision comprehensive analysis...")
            analysis_result = self.call_claude_advanced(vision_prompt, 4000, [img_bytes])
            
            if analysis_result and len(analysis_result.strip()) > 20:
                logger.info(f"Claude Vision analysis successful: {len(analysis_result)} characters")
                
                # Return structured result with both image bytes and analysis
                return {
                    "type": "image_analysis",
                    "filename": filename,
                    "analysis": analysis_result.strip(),
                    "image_bytes": img_bytes,
                    "size_kb": len(img_bytes) / 1024,
                    "vision_api_used": True,
                    "method": "claude_vision_primary"
                }
            else:
                logger.warning("Claude Vision returned empty or minimal result")
            
        except Exception as e:
            logger.warning(f"Claude Vision analysis failed: {e}")
        
        # FALLBACK CHAIN: OCR methods if vision analysis fails
        logger.info("Claude Vision failed, trying OCR fallbacks...")
        return self._fallback_ocr_processing(file_obj, filename)
    
    def _fallback_ocr_processing(self, file_obj, filename: str) -> Dict[str, Any]:
        """Fallback OCR processing if Claude Vision fails"""
        logger.info(f"Using OCR fallback methods for: {filename}")
        
        # Fallback 1: OCR.space API
        try:
            file_obj.seek(0)
            img_bytes = file_obj.read()
            
            logger.info("Attempting OCR.space API...")
            ocr_result = self._ocr_space_api(img_bytes)
            
            if ocr_result and len(ocr_result.strip()) > 10:
                logger.info(f"OCR.space API successful: {len(ocr_result)} characters")
                return {
                    "type": "ocr_text",
                    "filename": filename,
                    "text_content": ocr_result.strip(),
                    "image_bytes": img_bytes,
                    "ocr_method": "ocr_space_api",
                    "method": "ocr_space_fallback"
                }
            
        except Exception as e:
            logger.warning(f"OCR.space API failed: {e}")
        
        # Fallback 2: EasyOCR
        if self.easyocr_reader and PIL_AVAILABLE:
            try:
                file_obj.seek(0)
                img_bytes = file_obj.read()
                
                logger.info("Attempting EasyOCR...")
                import numpy as np
                
                # Convert to PIL Image
                image = Image.open(BytesIO(img_bytes))
                img_array = np.array(image)
                
                results = self.easyocr_reader.readtext(img_array)
                text_parts = [result[1] for result in results if result[2] > 0.3]  # confidence > 0.3
                
                if text_parts:
                    extracted_text = '\n'.join(text_parts)
                    logger.info(f"EasyOCR successful: {len(extracted_text)} characters")
                    return {
                        "type": "ocr_text",
                        "filename": filename,
                        "text_content": extracted_text,
                        "image_bytes": img_bytes,
                        "ocr_method": "easyocr",
                        "method": "easyocr_fallback"
                    }
                
            except Exception as e:
                logger.warning(f"EasyOCR failed: {e}")
        
        # Fallback 3: Tesseract
        if TESSERACT_AVAILABLE and PIL_AVAILABLE:
            try:
                file_obj.seek(0)
                img_bytes = file_obj.read()
                
                logger.info("Attempting Tesseract OCR...")
                
                image = Image.open(BytesIO(img_bytes))
                text = pytesseract.image_to_string(image)
                
                if text and len(text.strip()) > 10:
                    logger.info(f"Tesseract OCR successful: {len(text)} characters")
                    return {
                        "type": "ocr_text", 
                        "filename": filename,
                        "text_content": text.strip(),
                        "image_bytes": img_bytes,
                        "ocr_method": "tesseract",
                        "method": "tesseract_fallback"
                    }
                
            except Exception as e:
                logger.warning(f"Tesseract OCR failed: {e}")
        
        # Ultimate fallback - return error but keep image bytes for potential manual analysis
        file_obj.seek(0)
        img_bytes = file_obj.read()
        
        logger.error("All image processing methods failed")
        return {
            "type": "image_analysis_failed",
            "filename": filename,
            "error": "All vision and OCR methods failed",
            "image_bytes": img_bytes,
            "size_kb": len(img_bytes) / 1024,
            "method": "all_methods_failed"
        }

    def _ocr_space_api(self, img_bytes: bytes) -> str:
        """Use OCR.space free API for OCR"""
        try:
            files = {'file': ('image.jpg', img_bytes, 'image/jpeg')}
            data = {
                'apikey': OCR_SPACE_API_KEY,
                'language': 'eng',
                'isOverlayRequired': 'false',
                'detectOrientation': 'true',
                'scale': 'true',
                'OCREngine': '2',
                'isTable': 'true'
            }
            
            response = requests.post(OCR_SPACE_URL, files=files, data=data, timeout=30)
            response.raise_for_status()
            
            result = response.json()
            
            if result.get('IsErroredOnProcessing'):
                raise Exception(f"OCR.space error: {result.get('ErrorMessage', 'Unknown error')}")
            
            parsed_results = result.get('ParsedResults', [])
            if parsed_results:
                return parsed_results[0].get('ParsedText', '')
            
            return ''
            
        except Exception as e:
            logger.error(f"OCR.space API error: {e}")
            raise e

    def _process_word_file(self, file_obj) -> str:
        """Process Word documents"""
        try:
            file_obj.seek(0)
            
            # Try python-docx for .docx files
            if Document:
                try:
                    doc = Document(file_obj)
                    text = []
                    
                    for paragraph in doc.paragraphs:
                        text.append(paragraph.text)
                    
                    # Extract tables
                    for table in doc.tables:
                        for row in table.rows:
                            row_text = [cell.text for cell in row.cells]
                            text.append('\t'.join(row_text))
                    
                    return '\n'.join(text)
                    
                except Exception as docx_error:
                    logger.warning(f"python-docx failed: {docx_error}")
            else:
                logger.warning("python-docx not installed")
                
            # Fallback to textract if available
            if TEXTRACT_AVAILABLE:
                try:
                    file_obj.seek(0)
                    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
                        tmp.write(file_obj.read())
                        tmp.flush()
                        text = textract.process(tmp.name).decode('utf-8')
                        os.unlink(tmp.name)
                        return text
                except Exception as textract_error:
                    logger.warning(f"textract failed: {textract_error}")
            
            return "[Failed to extract text from Word document - missing dependencies]"
        
        except Exception as e:
            logger.error(f"Word processing failed: {e}")
            return ""

    def _process_powerpoint_file(self, file_obj) -> str:
        """Process PowerPoint presentations"""
        try:
            file_obj.seek(0)
            
            if not Presentation:
                logger.warning("python-pptx not installed")
                return "[PowerPoint processing requires python-pptx - install with: pip install python-pptx]"
            
            prs = Presentation(file_obj)
            text = []
            
            for slide in prs.slides:
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text.append(shape.text)
                
                if slide_text:
                    text.append("--- Slide ---")
                    text.extend(slide_text)
            
            return '\n'.join(text)
            
        except Exception as e:
            logger.error(f"PowerPoint processing failed: {e}")
            return ""

    def _process_sqlite_file(self, file_obj) -> Dict[str, pd.DataFrame]:
        """Process SQLite database files"""
        try:
            file_obj.seek(0)
            
            with tempfile.NamedTemporaryFile(suffix='.db', delete=False) as tmp:
                tmp.write(file_obj.read())
                tmp.flush()
                
                conn = sqlite3.connect(tmp.name)
                
                # Get all table names
                cursor = conn.execute("SELECT name FROM sqlite_master WHERE type='table';")
                tables = cursor.fetchall()
                
                result = {}
                for table_name, in tables:
                    try:
                        df = pd.read_sql_query(f"SELECT * FROM {table_name} LIMIT 10000", conn)
                        result[table_name] = self._clean_dataframe(df)
                    except Exception as e:
                        logger.warning(f"Failed to read table {table_name}: {e}")
                        continue
                
                conn.close()
                os.unlink(tmp.name)
                
                return result
                
        except Exception as e:
            logger.error(f"SQLite processing failed: {e}")
            return {}

    def _process_xml_html_file(self, file_obj) -> Any:
        """Process XML/HTML files"""
        try:
            file_obj.seek(0)
            content = file_obj.read()
            if isinstance(content, bytes):
                content = content.decode('utf-8', errors='ignore')
            
            # Try to extract structured data
            df = self.universal_data_extractor(content)
            if self._is_valid_dataframe(df):
                return df
            
            # Return raw content if no structure found
            return content
            
        except Exception as e:
            logger.error(f"XML/HTML processing failed: {e}")
            return ""

    def _process_zip_file(self, file_obj) -> Dict[str, Any]:
        """Process ZIP files - enhanced"""
        extracted = {}
        
        try:
            file_obj.seek(0)
            with zipfile.ZipFile(file_obj, 'r') as zip_ref:
                for file_info in zip_ref.filelist[:20]:  # Process up to 20 files
                    if not file_info.is_dir():
                        try:
                            with zip_ref.open(file_info) as extracted_file:
                                extracted_content = BytesIO(extracted_file.read())
                                
                                filename = file_info.filename
                                _, ext = os.path.splitext(filename.lower())
                                
                                if ext in ['.csv', '.tsv']:
                                    df = self._process_csv_file(extracted_content, ext)
                                    extracted[filename] = df
                                elif ext in ['.xlsx', '.xls']:
                                    df = self._process_excel_file(extracted_content)
                                    extracted[filename] = df
                                elif ext == '.json':
                                    data = self._process_json_file(extracted_content)
                                    extracted[filename] = data
                                elif ext in ['.txt', '.md']:
                                    text = self._process_text_file(extracted_content)
                                    extracted[filename] = text
                                elif ext == '.pdf':
                                    text = self._process_pdf_file(extracted_content)
                                    extracted[filename] = text
                                elif ext in ['.png', '.jpg', '.jpeg', '.gif', '.bmp']:
                                    text = self._process_image_file(extracted_content, filename)
                                    extracted[filename] = text
                                
                        except Exception as e:
                            logger.warning(f"Could not extract {file_info.filename}: {e}")
                            continue
        
        except Exception as e:
            logger.error(f"ZIP processing failed: {e}")
        
        return extracted

    def generate_question_specific_analysis_code(self, df: pd.DataFrame, question: str, sources: Dict, files: Dict = None) -> str:
        """Generate analysis code using LLM that follows the exact question format"""
        
        # Get data info
        data_info = {
            'shape': df.shape,
            'columns': list(df.columns)[:20],
            'dtypes': {k: str(v) for k, v in df.dtypes.to_dict().items()},
            'numeric_cols': df.select_dtypes(include=[np.number]).columns.tolist()[:10],
            'text_cols': df.select_dtypes(include=['object']).columns.tolist()[:10],
            'sample': df.head(3).to_dict() if len(df) > 0 else {}
        }
        
        # Enhanced prompt that GUARANTEES a result variable and handles df_clean properly
        analysis_prompt = f"""
You are an expert data analyst. Analyze the given data to answer the user's question EXACTLY as requested.

QUESTION: {question}

DATA INFO:
- Shape: {data_info['shape']}
- Columns: {data_info['columns']}
- Numeric columns: {data_info['numeric_cols']}
- Text columns: {data_info['text_cols']}
- Sample data: {str(data_info['sample'])[:500]}

CRITICAL REQUIREMENTS:
1. MUST start by creating df_clean = df.copy() and use df_clean throughout
2. MUST end your code with: result = [your_final_answer]
3. If question asks for JSON object, use: result = {{"key1": value1, "key2": value2}}
4. If question asks for JSON array, use: result = ["item1", "item2", "item3"]
5. For plots: ALWAYS use create_plot_base64() and include in result
6. Handle ALL edge cases and missing data with proper error handling
7. NEVER leave result undefined - always assign something to result
8. Use safe_column_access() for ALL DataFrame column access
9. Use convert_to_json_serializable() for all final results
10. Wrap all operations in try-except blocks

Generate Python code that analyzes the data and stores the final answer in 'result' variable:

```python
import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import seaborn as sns
import base64
import io
import re
import json
from datetime import datetime
import warnings
warnings.filterwarnings('ignore')
matplotlib.use('Agg')

def safe_column_access(df, col_name):
    \"\"\"Safely access columns by name, case-insensitive and partial matching\"\"\"
    if df is None or len(df) == 0:
        return pd.Series()
    
    # Exact match first
    if col_name in df.columns:
        return df[col_name]
    
    # Case insensitive match
    for col in df.columns:
        if col.lower() == col_name.lower():
            return df[col]
    
    # Partial match
    for col in df.columns:
        if col_name.lower() in col.lower() or col.lower() in col_name.lower():
            return df[col]
    
    # Return first column as fallback
    return df.iloc[:, 0] if len(df.columns) > 0 else pd.Series()

def safe_numeric_convert(series):
    \"\"\"Safely extract numbers from text\"\"\"
    try:
        if series.dtype == 'object':
            # Extract numeric parts and convert
            numeric_series = pd.to_numeric(
                series.astype(str).str.replace(r'[^0-9.-]', '', regex=True), 
                errors='coerce'
            )
            return numeric_series
        return pd.to_numeric(series, errors='coerce')
    except Exception as e:
        return pd.to_numeric(series, errors='coerce')

def create_plot_base64(fig=None, max_size_kb=100):
    \"\"\"Create base64 plot - GUARANTEED to return a valid data URI\"\"\"
    try:
        if fig is None:
            fig = plt.gcf()
        
        # Try different DPI values to control size
        for dpi in [72, 60, 50, 40, 30]:
            try:
                buf = io.BytesIO()
                fig.savefig(buf, format='png', dpi=dpi, bbox_inches='tight', 
                           facecolor='white', edgecolor='none', pad_inches=0.1)
                buf.seek(0)
                
                img_data = base64.b64encode(buf.read()).decode('utf-8')
                buf.close()
                
                # Check size
                size_kb = len(img_data) * 3 / 4 / 1024
                if size_kb <= max_size_kb:
                    plt.close(fig)
                    return "data:image/png;base64," + img_data
            except Exception as inner_e:
                continue
        
        # If all attempts fail, return minimal plot
        plt.close(fig)
        return "data:image/png;base64,iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJAAAADUlEQVR42mNkYPhfDwAChwGA60e6kgAAAABJRU5ErkJggg=="
        
    except Exception as e:
        return "data:image/png;base64,iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJAAAADUlEQVR42mNkYPhfDwAChwGA60e6kgAAAABJRU5ErkJggg=="

def convert_to_json_serializable(obj):
    \"\"\"Convert numpy/pandas types to JSON serializable types\"\"\"
    try:
        if obj is None or pd.isna(obj):
            return None
        elif isinstance(obj, (np.integer, np.int64, np.int32)):
            return int(obj)
        elif isinstance(obj, (np.floating, np.float64, np.float32)):
            if np.isnan(obj):
                return None
            return float(obj)
        elif isinstance(obj, np.ndarray):
            return [convert_to_json_serializable(item) for item in obj.tolist()]
        elif isinstance(obj, pd.Series):
            return [convert_to_json_serializable(item) for item in obj.tolist()]
        elif isinstance(obj, pd.Timestamp):
            return obj.isoformat()
        elif isinstance(obj, (list, tuple)):
            return [convert_to_json_serializable(item) for item in obj]
        elif isinstance(obj, dict):
            return {{k: convert_to_json_serializable(v) for k, v in obj.items()}}
        elif hasattr(obj, 'item'):
            return convert_to_json_serializable(obj.item())
        else:
            return obj
    except Exception as e:
        return str(obj) if obj is not None else None

# Initialize result to ensure it always exists
result = None

try:
    # STEP 1: Create clean DataFrame copy
    df_clean = df.copy()
    
    if df_clean is None or len(df_clean) == 0:
        result = {{"error": "No data available for analysis", "question": "{question[:200]}"}}
    else:
        # STEP 2: Data preprocessing and cleaning
        # Remove completely empty rows/columns
        df_clean = df_clean.dropna(how='all').dropna(axis=1, how='all')
        
        # Clean column names
        df_clean.columns = [str(col).strip() for col in df_clean.columns]
        
        # Convert problematic string numbers to numeric where possible
        for col in df_clean.select_dtypes(include=['object']).columns:
            try:
                # Try to convert if it looks numeric
                sample = df_clean[col].dropna().head(10)
                if len(sample) > 0:
                    # Check if most values look numeric
                    numeric_count = 0
                    for val in sample:
                        if re.match(r'^[+-]?(\d+\.?\d*|\.\d+), str(val).strip()):
                            numeric_count += 1
                    
                    if numeric_count / len(sample) > 0.7:  # 70% look numeric
                        df_clean[col + '_numeric'] = safe_numeric_convert(df_clean[col])
            except Exception as col_error:
                pass  # Skip problematic columns
        
        # STEP 3: MAIN ANALYSIS - Generate based on question
        # [The LLM will fill this section with specific analysis]
        
        # STEP 4: Ensure result is properly formatted and JSON serializable
        if result is None:
            result = {{
                "error": "Analysis completed but no specific result generated",
                "data_shape": df_clean.shape,
                "columns": list(df_clean.columns)[:10],
                "question": "{question[:200]}"
            }}

except Exception as e:
    # STEP 5: Fallback result in case of any error
    result = {{
        "error": f"Analysis failed: {{str(e)}}",
        "data_shape": df.shape if 'df' in locals() else "unknown",
        "question": "{question[:200]}",
        "error_type": type(e).__name__
    }}

# STEP 6: Final conversion to ensure JSON compatibility
try:
    if isinstance(result, dict):
        result = {{k: convert_to_json_serializable(v) for k, v in result.items()}}
    elif isinstance(result, list):
        result = [convert_to_json_serializable(item) for item in result]
    else:
        result = convert_to_json_serializable(result)
except Exception as final_error:
    result = {{
        "error": "Result conversion failed", 
        "raw_result": str(result)[:500],
        "conversion_error": str(final_error)
    }}

# GUARANTEE: result variable always exists and is JSON serializable
assert result is not None, "Result must not be None"
```

IMPORTANT: Your code MUST:
1. Always create df_clean = df.copy() at the start
2. Use df_clean for ALL operations 
3. End with a valid 'result' variable assignment
4. Handle the "T2257844554" type errors by using safe_numeric_convert()
5. Use proper array operations with .any() or .all() for boolean arrays
6. Never use ambiguous truth values on arrays
"""

        try:
            response = self.call_claude_advanced(analysis_prompt, 6000)
            
            # Extract code from response
            if '```python' in response:
                code = response.split('```python')[1].split('```')[0]
            elif '```' in response:
                code = response.split('```')[1].split('```')[0]
            else:
                code = response
            
            # GUARANTEE that code has df_clean definition and result assignment
            code = code.strip()
            
            #Ensure df_clean is defined at the start of the code
            if 'df_clean' not in code or not re.search(r'df_clean\s*=\s*df\.copy\(\)', code):
                df_clean_line = 'df_clean = df.copy()'
                code = df_clean_line + '\n' + code
            
            # Ensure result assignment exists
            if 'result =' not in code:
                code += f"""

# Emergency fallback - ensure result exists
if 'result' not in locals() or result is None:
    result = {{
        "error": "No analysis result generated",
        "question": "{question[:200]}",
        "data_available": True,
        "data_shape": df_clean.shape if 'df_clean' in locals() else df.shape
    }}
"""
            
            return code
            
        except Exception as e:
            logger.error(f"Code generation failed: {e}")
            return self._guaranteed_fallback_analysis_code(df, question)

    def _guaranteed_fallback_analysis_code(self, df: pd.DataFrame, question: str) -> str:
        """Guaranteed fallback that ALWAYS produces a result"""
        return f"""
# GUARANTEE fallback analysis - ALWAYS produces result
import pandas as pd
import numpy as np
import json

def safe_column_access(df, col_name):
    if df is None or len(df) == 0:
        return pd.Series()
    if col_name in df.columns:
        return df[col_name]
    for col in df.columns:
        if col.lower() == col_name.lower():
            return df[col]
    for col in df.columns:
        if col_name.lower() in col.lower():
            return df[col]
    return df.iloc[:, 0] if len(df.columns) > 0 else pd.Series()

def convert_to_json_serializable(obj):
    try:
        if obj is None or pd.isna(obj):
            return None
        elif isinstance(obj, (np.integer, np.int64, np.int32)):
            return int(obj)
        elif isinstance(obj, (np.floating, np.float64, np.float32)):
            if np.isnan(obj):
                return None
            return float(obj)
        elif isinstance(obj, np.ndarray):
            return obj.tolist()
        elif isinstance(obj, pd.Series):
            return obj.tolist()
        elif isinstance(obj, pd.Timestamp):
            return obj.isoformat()
        elif isinstance(obj, (list, tuple)):
            return [convert_to_json_serializable(item) for item in obj]
        elif isinstance(obj, dict):
            return {{k: convert_to_json_serializable(v) for k, v in obj.items()}}
        else:
            return obj
    except:
        return str(obj) if obj is not None else None

# Create clean DataFrame
df_clean = df.copy() if df is not None else pd.DataFrame()

# Determine output format from question
question_lower = "{question.lower()}"

try:
    if df_clean is not None and len(df_clean) > 0:
        if "json array" in question_lower or "respond with a json array" in question_lower:
            # Return array format with basic data info
            result = [
                f"Data shape: {{df_clean.shape}}",
                f"Columns: {{list(df_clean.columns)[:5]}}",
                "Fallback analysis - LLM code generation failed",
                "Please check data format and retry"
            ]
        else:
            # Return object format with data summary
            result = {{
                "error": "LLM code generation failed - using guaranteed fallback",
                "data_shape": convert_to_json_serializable(df_clean.shape),
                "columns": convert_to_json_serializable(list(df_clean.columns)[:10]),
                "sample_data": convert_to_json_serializable(df_clean.head(2).to_dict()) if len(df_clean) > 0 else {{}},
                "question": "{question[:200]}",
                "status": "guaranteed_fallback_used",
                "numeric_columns": convert_to_json_serializable(df_clean.select_dtypes(include=[np.number]).columns.tolist()[:5]),
                "text_columns": convert_to_json_serializable(df_clean.select_dtypes(include=['object']).columns.tolist()[:5])
            }}
    else:
        result = {{
            "error": "No valid data available",
            "question": "{question[:200]}",
            "data_available": False
        }}
    
    # Ensure JSON compatibility
    if isinstance(result, dict):
        result = {{k: convert_to_json_serializable(v) for k, v in result.items()}}
    elif isinstance(result, list):
        result = [convert_to_json_serializable(item) for item in result]
        
except Exception as e:
    # Ultimate fallback
    result = {{
        "error": f"Even guaranteed fallback failed: {{str(e)}}",
        "question": "{question[:100]}",
        "data_available": df is not None and len(df) > 0 if df is not None else False
    }}

# GUARANTEE: result variable always exists
assert result is not None, "Result must not be None"
"""

    def execute_analysis_safely(self, code: str, df: pd.DataFrame, question: str, sources: Dict, files: Dict = None) -> Any:
        """Execute analysis code with enhanced error handling and LLM-based fixing"""
        
        for attempt in range(1, self.max_retries + 1):
            try:
                logger.info(f"Executing analysis attempt {attempt}/{self.max_retries}")
                
                # Set up secure execution environment
                exec_globals = {
                    'pd': pd, 'np': np, 'plt': plt, 'sns': sns,
                    'json': json, 'base64': base64, 'io': io,
                    'df': df.copy(), 'question': question, 'sources': sources,
                    'datetime': datetime, 're': re, 'warnings': warnings,
                    'matplotlib': matplotlib
                }
                
                # Add helper functions to execution environment
                def safe_column_access(df, col_name):
                    """Safely access columns"""
                    if col_name in df.columns:
                        return df[col_name]
                    for col in df.columns:
                        if col.lower() == col_name.lower():
                            return df[col]
                    for col in df.columns:
                        if col_name.lower() in col.lower():
                            return df[col]
                    return df.iloc[:, 0] if len(df.columns) > 0 else pd.Series()

                def safe_numeric_convert(series):
                    """Safely extract numbers from text"""
                    if series.dtype == 'object':
                        numeric = pd.to_numeric(
                            series.astype(str).str.extract(r'([+-]?\\d*[.,]?\\d+)')[0].str.replace(',', ''), 
                            errors='coerce'
                        )
                        return numeric
                    return series

                def create_plot_base64(fig=None, max_size_kb=100):
                    """Create base64 encoded plot - optimized for size"""
                    if fig is None:
                        fig = plt.gcf()
                    
                    img_str = None  # Initialize to avoid undefined variable error
                    
                    # Try different DPI values to stay under size limit
                    for dpi in [72, 60, 50, 40]:
                        buf = io.BytesIO()
                        try:
                            fig.savefig(buf, format='png', dpi=dpi, bbox_inches='tight', 
                                        facecolor='white', edgecolor='none', pad_inches=0.1)
                            buf.seek(0)
                            img_str = base64.b64encode(buf.read()).decode('utf-8')
                            buf.close()
                            
                            # Check size (approximate)
                            size_kb = len(img_str) * 3 / 4 / 1024
                            if size_kb <= max_size_kb:
                                plt.close(fig)
                                return f"data:image/png;base64,{img_str}"
                        except Exception as e:
                            buf.close()
                            continue
                    
                    # Return with the last successful encoding, or a fallback
                    plt.close(fig)
                    if img_str:
                        return f"data:image/png;base64,{img_str}"
                    else:
                        # Fallback: create a simple plot
                        return "data:image/png;base64,iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJAAAADUlEQVR42mNkYPhfDwAChwGA60e6kgAAAABJRU5ErkJggg=="

                def convert_to_json_serializable(obj):
                    """Convert numpy/pandas types to JSON serializable types"""
                    if isinstance(obj, (np.integer, np.int64, np.int32)):
                        return int(obj)
                    elif isinstance(obj, (np.floating, np.float64, np.float32)):
                        return float(obj)
                    elif isinstance(obj, np.ndarray):
                        return obj.tolist()
                    elif isinstance(obj, pd.Series):
                        return obj.tolist()
                    elif pd.isna(obj) or obj is None:
                        return None
                    elif isinstance(obj, pd.Timestamp):
                        return obj.isoformat()
                    elif hasattr(obj, 'item'):  # numpy scalars
                        return obj.item()
                    return obj

                # Add functions to execution environment
                exec_globals.update({
                    'safe_column_access': safe_column_access,
                    'safe_numeric_convert': safe_numeric_convert,
                    'create_plot_base64': create_plot_base64,
                    'convert_to_json_serializable': convert_to_json_serializable
                })
                
                if files:
                    exec_globals.update({f'file_{k.replace(".", "_")}': v for k, v in files.items()})
                
                # Clear any existing plots
                plt.clf()
                plt.close('all')
                
                # Execute with timeout and better error handling
                with timeout_context(90):  # 90 second timeout for execution
                    try:
                        exec(code, exec_globals)
                    except NameError as ne:
                        logger.error(f"NameError in generated code: {ne}")
                        # Try to add missing common variables
                        if 'img_data' in str(ne):
                            exec_globals['img_data'] = None
                        elif 'img_str' in str(ne):
                            exec_globals['img_str'] = None
                        # Re-execute with fixes
                        exec(code, exec_globals)
                    except Exception as ee:
                        logger.error(f"Execution error: {ee}")
                        raise ee
                
                # Get result
                if 'result' in exec_globals:
                    result = exec_globals['result']
                    logger.info(f"Analysis successful: {type(result)}")
                    
                    # Recursively convert numpy/pandas types to JSON serializable
                    def deep_convert(obj):
                        if isinstance(obj, dict):
                            return {k: deep_convert(v) for k, v in obj.items()}
                        elif isinstance(obj, list):
                            return [deep_convert(item) for item in obj]
                        elif isinstance(obj, tuple):
                            return tuple(deep_convert(item) for item in obj)
                        else:
                            return convert_to_json_serializable(obj)
                    
                    result = deep_convert(result)
                    
                    # Ensure base64 images have proper prefix
                    def ensure_base64_prefix(obj):
                        if isinstance(obj, str):
                            # Check if it's a base64 image without prefix
                            if (obj.startswith('iVBORw0KGgo') or obj.startswith('/9j/') or 
                                obj.startswith('R0lGODlh') or obj.startswith('UklGRg==') or
                                len(obj) > 50 and obj.replace('+', '').replace('/', '').replace('=', '').isalnum()):
                                return f"data:image/png;base64,{obj}"
                            return obj
                        elif isinstance(obj, dict):
                            return {k: ensure_base64_prefix(v) for k, v in obj.items()}
                        elif isinstance(obj, list):
                            return [ensure_base64_prefix(item) for item in obj]
                        else:
                            return obj
                    
                    result = ensure_base64_prefix(result)
                    return result
                else:
                    logger.warning("No 'result' variable found in code execution")
                    return {"error": "No result produced"}
                    
            except Exception as e:
                logger.warning(f"Analysis attempt {attempt} failed: {e}")
                
                if attempt == self.max_retries:
                    logger.error(f"All analysis attempts failed. Last error: {e}")
                    return {"error": f"Analysis failed: {str(e)}"}
                
                # Try to fix the code using LLM for next attempt
                try:
                    if self.time_remaining() > 30:
                        code = self._fix_analysis_code_with_llm(code, str(e), df, question)
                except:
                    pass  # If fixing fails, just retry with same code
        
        return {"error": "Analysis execution failed after all retries", "question": question[:200]}

    def _fix_analysis_code_with_llm(self, code: str, error: str, df: pd.DataFrame, question: str) -> str:
        """Use LLM to fix common code errors"""
        
        fix_prompt = f"""
Fix this Python code that failed with error: {error}

Original code:
```python
{code[:2000]}
```

Error details: {error}

Common fixes needed:
1. For "could not convert string to float" errors: Use safe_numeric_convert() for ALL numeric operations
2. For "truth value of array" errors: Use len(df) > 0 instead of if df, use .any() or .all() for boolean arrays
3. Use safe_column_access(df, 'column_name') for ALL column access
4. Use convert_to_json_serializable() for all final results
5. Use create_plot_base64() for all plots with proper size limits
6. Always create df_clean = df.copy() at the start
7. Handle NaN/None values properly with pd.isna() checks
8. Ensure proper data type conversions
9. Add try-catch blocks for error handling
10. Never use ambiguous boolean operations on DataFrames

Question context: {question[:200]}

Return the complete fixed code that will definitely work:
"""
        
        try:
            response = self.call_claude_advanced(fix_prompt, 3000)
            
            # Extract code from response
            if '```python' in response:
                fixed_code = response.split('```python')[1].split('```')[0]
            elif '```' in response:
                fixed_code = response.split('```')[1].split('```')[0]
            else:
                fixed_code = response
            
            # Ensure df_clean is defined
            if 'df_clean = df.copy()' not in fixed_code:
                fixed_code = 'df_clean = df.copy()\n' + fixed_code
            
            logger.info("LLM successfully generated fixed code")
            return fixed_code.strip()
        except Exception as e:
            logger.warning(f"LLM code fixing failed: {e}")
            return code  # Return original if fixing fails

    async def process_question_universally(self, question: str, files: Dict = None) -> Any:
        """Main processing pipeline - enhanced with better LLM integration"""
        
        self.start_time = time.time()
        
        try:
            logger.info(f"=== Processing question universally ===")
            logger.info(f"Question: {question[:200]}...")
            
            # Step 1: Detect all possible data sources
            sources = self.detect_data_sources(question, files)
            logger.info(f"Detected sources: {sources}")
            
            # Step 2: Process attached files first (highest priority)
            processed_files = {}
            primary_df = None
            
            if files:
                logger.info("Processing attached files...")
                processed_files = self.process_attached_files(files)
                
                # Find the best DataFrame to work with
                for filename, data in processed_files.items():
                    logger.info(f"Checking processed file {filename}: type={type(data)}")
                    if self._is_valid_dataframe(data):
                        primary_df = data
                        logger.info(f"Using {filename} as primary data source: {data.shape}")
                        break
                    elif isinstance(data, dict) and 'error' not in data:
                        # Try to convert dict to DataFrame if structured
                        try:
                            if isinstance(list(data.values())[0], list):
                                df = pd.DataFrame(data)
                                if self._is_valid_dataframe(df):
                                    primary_df = df
                                    logger.info(f"Converted {filename} dict to DataFrame: {df.shape}")
                                    break
                        except:
                            pass
                    elif isinstance(data, str) and len(data) > 100:
                        # Try to extract structured data from text (like OCR results)
                        try:
                            extracted_df = self.universal_data_extractor(data)
                            if self._is_valid_dataframe(extracted_df):
                                primary_df = extracted_df
                                logger.info(f"Extracted DataFrame from text in {filename}: {extracted_df.shape}")
                                break
                        except:
                            pass
                    elif isinstance(data, dict) and data.get('type') == 'image_analysis' and data.get('analysis'):
                        # Handle Claude Vision analysis results
                        logger.info(f"Found Claude Vision analysis for {filename}")
                        # Try to extract structured data from the analysis
                        try:
                            analysis_text = data['analysis']
                            extracted_df = self.universal_data_extractor(analysis_text)
                            if self._is_valid_dataframe(extracted_df):
                                primary_df = extracted_df
                                logger.info(f"Extracted DataFrame from Vision analysis: {extracted_df.shape}")
                                break
                            else:
                                # If no tabular data found, we'll use the image analysis directly
                                processed_files[f"{filename}_vision_analysis"] = {
                                    "type": "vision_analysis",
                                    "content": analysis_text,
                                    "source_file": filename
                                }
                        except:
                            pass
                
                if primary_df is None:
                    logger.warning(f"No valid DataFrame found in processed files")
                    for filename, data in processed_files.items():
                        logger.warning(f"  {filename}: {type(data)}")
                else:
                    logger.info(f"Found primary DataFrame: {primary_df.shape}")
            
            # Step 3: Web scraping if URLs found and no primary data yet
            if self._is_empty_or_none(primary_df) and sources.get('urls') and self.time_remaining() > 50:
                logger.info("Attempting web scraping...")
                for url in sources['urls'][:2]:
                    try:
                        if self.time_remaining() < 40:
                            break
                        
                        html_content = await self.universal_web_scraper(url)
                        df = self.universal_data_extractor(html_content, url)
                        
                        if self._is_valid_dataframe(df):
                            primary_df = df
                            logger.info(f"Scraped data from {url}: {df.shape}")
                            break
                    except Exception as e:
                        logger.warning(f"Scraping failed for {url}: {e}")
                        continue
            
            # Step 4: Database/file system queries if no primary data yet
            if self._is_empty_or_none(primary_df) and (sources.get('s3_paths') or sources.get('database_urls') or sources.get('query_hints')) and self.time_remaining() > 40:
                logger.info("Attempting database/file queries...")
                db_result = self.universal_database_handler(question, sources)
                
                if db_result and not (isinstance(db_result, dict) and 'error' in db_result):
                    # Convert result to DataFrame if possible
                    if isinstance(db_result, list) and len(db_result) > 0:
                        if isinstance(db_result[0], dict):
                            primary_df = pd.DataFrame(db_result)
                        else:
                            primary_df = pd.DataFrame(db_result, columns=[f'col_{i}' for i in range(len(db_result[0]))])
                    else:
                        # Return direct result for simple queries
                        def make_json_serializable(obj):
                            if isinstance(obj, (np.integer, np.int64, np.int32)):
                                return int(obj)
                            elif isinstance(obj, (np.floating, np.float64, np.float32)):
                                return float(obj)
                            elif isinstance(obj, np.ndarray):
                                return obj.tolist()
                            elif isinstance(obj, list):
                                return [make_json_serializable(item) for item in obj]
                            elif isinstance(obj, dict):
                                return {k: make_json_serializable(v) for k, v in obj.items()}
                            return obj
                        
                        return make_json_serializable(db_result)
            
            # Step 5: If we have data, analyze it using LLM-generated code
            if self._is_valid_dataframe(primary_df):
                logger.info(f"Analyzing data: {primary_df.shape}")
                
                # Generate analysis code that follows the question format exactly using LLM
                analysis_code = self.generate_question_specific_analysis_code(
                    primary_df, question, sources, processed_files
                )
                
                # Execute analysis with enhanced error handling
                result = self.execute_analysis_safely(
                    analysis_code, primary_df, question, sources, processed_files
                )
                
                return result
            
            # Step 6: If we have non-tabular data (like text from images/PDFs or vision analysis), use LLM to analyze it directly
            if processed_files and not primary_df:
                logger.info("No tabular data found, checking for text/vision data to analyze...")
                
                # Check for Claude Vision analysis first
                vision_analyses = []
                text_contents = []
                
                for filename, data in processed_files.items():
                    if isinstance(data, dict):
                        if data.get('type') == 'image_analysis' and data.get('analysis'):
                            vision_analyses.append(f"=== Image Analysis: {filename} ===\n{data['analysis']}")
                        elif data.get('type') == 'vision_analysis':
                            vision_analyses.append(f"=== Vision Analysis: {data.get('source_file', filename)} ===\n{data['content']}")
                        elif data.get('type') == 'ocr_text' and data.get('text_content'):
                            text_contents.append(f"=== OCR Text from {filename} ===\n{data['text_content']}")
                    elif isinstance(data, str) and len(data.strip()) > 50:
                        text_contents.append(f"=== {filename} ===\n{data}")
                
                # Prioritize vision analysis over OCR text
                all_content = vision_analyses + text_contents
                
                if all_content and self.time_remaining() > 25:
                    combined_content = "\n\n".join(all_content)
                    
                    # Use LLM to analyze the content directly, with image context
                    analysis_prompt = f"""
You are analyzing content extracted from files, including comprehensive image analysis. 

USER QUESTION: {question}

EXTRACTED CONTENT:
{combined_content[:4000]}

INSTRUCTIONS:
1. Answer the user's question based on the extracted content
2. If the question asks for specific format (JSON array/object), follow it exactly
3. Look for patterns, data, numbers, lists, tables, or structured information
4. For image analysis, focus on visual insights, text extraction, and data interpretation
5. If you find tabular data in the content, structure it appropriately
6. Provide comprehensive insights based on what you can see and extract
7. If there are charts, graphs, or visual data, describe and interpret them

Provide a detailed answer that directly addresses: "{question}"
"""
                    
                    try:
                        response = self.call_claude_advanced(analysis_prompt, 4000)
                        
                        # Try to parse as JSON if it looks like JSON
                        if response.strip().startswith('{') or response.strip().startswith('['):
                            try:
                                return json.loads(response)
                            except:
                                pass
                        
                        # Return structured response
                        return {
                            "analysis": response,
                            "source": "comprehensive_content_analysis",
                            "files_processed": list(processed_files.keys()),
                            "question": question[:200],
                            "content_types_found": [
                                "vision_analysis" if vision_analyses else None,
                                "ocr_text" if text_contents else None,
                                "other_text" if any(isinstance(data, str) for data in processed_files.values()) else None
                            ]
                        }
                        
                    except Exception as e:
                        logger.warning(f"Content analysis failed: {e}")
                
                # If we have image analysis but couldn't process it above, return it directly
                for filename, data in processed_files.items():
                    if isinstance(data, dict) and data.get('type') == 'image_analysis' and data.get('analysis'):
                        return {
                            "image_analysis": data['analysis'],
                            "filename": filename,
                            "source": "claude_vision_direct",
                            "question": question[:200],
                            "image_size_kb": data.get('size_kb', 0)
                        }
            
            # Step 7: No data found - use LLM to provide helpful response
            elapsed = time.time() - self.start_time
            logger.warning(f"No processable data found after {elapsed:.2f} seconds")
            
            # Use LLM to generate helpful response based on question
            if self.time_remaining() > 20:
                try:
                    no_data_prompt = f"""
The user asked: {question}

No data sources were found or could be processed. Generate a helpful response that:
1. Acknowledges the question
2. Explains what data would be needed
3. Suggests how to provide the required data
4. If the question expects a specific format (JSON array/object), return that format with explanatory messages

Available file info: {list(processed_files.keys()) if processed_files else "No files"}
Detected sources: {sources}

Return a JSON response appropriate for the question format.
"""
                    
                    llm_response = self.call_claude_advanced(no_data_prompt, 1000)
                    
                    # Try to parse as JSON if it looks like JSON
                    if llm_response.strip().startswith('{') or llm_response.strip().startswith('['):
                        try:
                            return json.loads(llm_response)
                        except:
                            pass
                    
                    # Return structured response
                    return {
                        "error": "No data sources could be processed",
                        "question": question[:200],
                        "suggestion": llm_response,
                        "sources_detected": sources,
                        "files_processed": list(processed_files.keys()) if processed_files else [],
                        "processing_time": elapsed
                    }
                    
                except Exception as e:
                    logger.warning(f"LLM no-data response failed: {e}")
            
            return {
                "error": "No data sources could be processed",
                "sources_detected": sources,
                "files_processed": list(processed_files.keys()) if processed_files else [],
                "processing_time": elapsed,
                "suggestion": "Please provide data files (CSV, Excel, JSON, images with data, PDFs), URLs with structured data, or database connections"
            }
            
        except Exception as e:
            elapsed = time.time() - self.start_time if self.start_time else 0
            logger.error(f"Universal processing failed after {elapsed:.2f} seconds: {e}")
            
            return {
                "error": f"Processing failed: {str(e)}",
                "processing_time": elapsed,
                "error_type": type(e).__name__
            }

# Initialize the universal agent
agent = None

@app.route('/api/', methods=['POST'])
async def analyze_data():
    """Universal data analysis API endpoint"""
    global agent
    
    # Initialize agent if needed
    if agent is None:
        if not os.getenv('ANTHROPIC_API_KEY'):
            return jsonify({"error": "ANTHROPIC_API_KEY environment variable not set"}), 500
        agent = UniversalDataAnalystAgent()
    
    try:
        # Extract question and files from request
        question = ""
        attached_files = {}
        
        logger.info(f"Request method: {request.method}")
        logger.info(f"Content type: {request.content_type}")
        logger.info(f"Files: {list(request.files.keys()) if request.files else 'None'}")
        logger.info(f"Form: {list(request.form.keys()) if request.form else 'None'}")
        
        # Handle file uploads
        if request.files:
            for key, file in request.files.items():
                logger.info(f"Processing upload: {key} -> {file.filename}")
                
                # Check if this is the question file
                if key == 'questions.txt' or (file.filename and 'question' in file.filename.lower()):
                    question = file.read().decode('utf-8').strip()
                    logger.info(f"Found question: {question[:100]}...")
                else:
                    # Store as attached file for analysis
                    file.seek(0)  # Reset file pointer
                    attached_files[file.filename or key] = file
        
        # Fallback methods for getting question
        if not question:
            # Try form data
            if request.form:
                for key, value in request.form.items():
                    if 'question' in key.lower() or len(value.strip()) > 10:
                        question = value.strip()
                        break
            
            # Try JSON body
            elif request.is_json:
                json_data = request.get_json()
                question = json_data.get('question', '') if json_data else ''
            
            # Try raw body
            else:
                raw_data = request.get_data(as_text=True)
                if raw_data and len(raw_data.strip()) > 5:
                    question = raw_data.strip()
        
        # Validate question
        if not question or len(question.strip()) < 5:
            return jsonify({
                "error": "No valid question provided",
                "details": "Expected questions.txt file or question in request body",
                "received_files": list(attached_files.keys()),
                "help": "Upload a questions.txt file with your question, or include the question in the request body"
            }), 400
        
        logger.info(f"Processing question: {question[:200]}...")
        logger.info(f"Attached files: {list(attached_files.keys())}")
        
        # Process the question universally using LLM-driven analysis
        result = await agent.process_question_universally(question, attached_files)
        
        # Format and return response
        if isinstance(result, dict) and result.get('error') and 'processing_time' not in result:
            logger.error(f"Processing error: {result}")
            return jsonify(result), 500
        
        # Use custom JSON encoder to handle numpy/pandas types
        try:
            response = app.response_class(
                response=json.dumps(result, cls=NumpyJSONEncoder, ensure_ascii=False),
                status=200,
                mimetype='application/json'
            )
            return response
                
        except Exception as json_error:
            logger.error(f"JSON serialization error: {json_error}")
            # Fallback: convert everything to strings
            def stringify_all(obj):
                if isinstance(obj, dict):
                    return {k: stringify_all(v) for k, v in obj.items()}
                elif isinstance(obj, list):
                    return [stringify_all(item) for item in obj]
                elif isinstance(obj, (np.integer, np.floating, np.ndarray, pd.Series)):
                    return str(obj)
                else:
                    return obj
            
            safe_result = stringify_all(result)
            return jsonify({
                "result": safe_result,
                "warning": "Some values converted to strings due to serialization issues"
            })
        
    except Exception as e:
        logger.error(f"API endpoint error: {e}")
        traceback.print_exc()
        return jsonify({
            "error": f"Request processing failed: {str(e)}",
            "type": type(e).__name__,
            "help": "Please check your request format and try again"
        }), 500

@app.route('/health', methods=['GET'])
def health_check():
    """Health check endpoint"""
    return jsonify({
        "status": "healthy",
        "timestamp": datetime.now().isoformat(),
        "version": "universal-data-analyst-5.0-comprehensive",
        "capabilities": [
            "llm_driven_analysis", "web_scraping", "comprehensive_file_processing", 
            "database_queries", "data_analysis", "visualization", 
            "multi_format_support", "json_serialization", 
            "question_specific_analysis", "error_recovery",
            "adaptive_code_generation", "dynamic_response_formatting",
            "hybrid_ocr", "claude_vision", "ocr_space_api", "easyocr_fallback",
            "image_processing", "pdf_extraction", "word_excel_processing",
            "powerpoint_processing", "sqlite_processing", "zip_archive_processing",
            "enhanced_error_handling", "robust_data_cleaning"
        ],
        "supported_file_types": {
            "tabular": ["csv", "tsv", "xlsx", "xls", "parquet", "json", "xml", "html"],
            "images": ["png", "jpg", "jpeg", "gif", "bmp", "tiff", "webp", "svg"],
            "documents": ["pdf", "doc", "docx", "txt", "md", "rtf"],
            "presentations": ["ppt", "pptx"],
            "databases": ["db", "sqlite", "sqlite3"],
            "archives": ["zip", "tar", "gz", "7z", "rar"],
            "other": ["log", "conf", "ini", "yaml", "yml"]
        },
        "optional_dependencies": {
            "playwright": PLAYWRIGHT_AVAILABLE,
            "aiohttp": AIOHTTP_AVAILABLE,
            "easyocr": EASYOCR_AVAILABLE,
            "pymupdf": PYMUPDF_AVAILABLE,
            "tesseract": TESSERACT_AVAILABLE,
            "pillow": PIL_AVAILABLE,
            "textract": TEXTRACT_AVAILABLE,
            "aws": AWS_AVAILABLE,
            "bigquery": BIGQUERY_AVAILABLE,
            "psycopg2": PSYCOPG2_AVAILABLE,
            "mysql": MYSQL_AVAILABLE,
            "mongodb": MONGODB_AVAILABLE,
            "sqlalchemy": SQLALCHEMY_AVAILABLE,
            "openpyxl": OPENPYXL_AVAILABLE,
            "asyncpg": ASYNCPG_AVAILABLE,
            "duckdb": DUCKDB_AVAILABLE
        }
    })

@app.route('/capabilities', methods=['GET'])
def get_capabilities():
    """Return supported capabilities"""
    return jsonify({
        "supported_formats": {
            "tabular": [".csv", ".tsv", ".xlsx", ".xls", ".parquet", ".json", ".xml", ".html"],
            "images": [".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff", ".webp", ".svg"],
            "documents": [".pdf", ".doc", ".docx", ".txt", ".md", ".rtf"],
            "presentations": [".ppt", ".pptx"],
            "databases": [".db", ".sqlite", ".sqlite3"],
            "archives": [".zip", ".tar", ".gz", ".7z"],
            "other": [".log", ".conf", ".ini", ".yaml", ".yml"]
        },
        "data_sources": [
            "HTTP/HTTPS URLs", "FTP URLs", "S3 paths", "Google Cloud Storage",
            "PostgreSQL", "MySQL", "MongoDB", "SQLite", "DuckDB",
            "API endpoints", "File uploads", "Zip archives", "Images with OCR",
            "PDF documents", "Word documents", "PowerPoint presentations"
        ],
        "ocr_capabilities": [
            "Claude Vision API (primary)", "OCR.space API (free tier)",
            "EasyOCR (offline backup)", "Tesseract (fallback)"
        ],
        "analysis_types": [
            "Statistical analysis", "Data visualization", "Correlation analysis",
            "Regression analysis", "Time series analysis", "Text analysis",
            "Image OCR analysis", "PDF content extraction", "Web scraping",
            "Custom question-specific analysis", "Multi-format data integration"
        ],
        "output_formats": [
            "JSON arrays", "JSON objects", "Base64 images", "Plain text",
            "Statistical metrics", "Structured data", "Custom formats",
            "Dynamic format adaptation based on question"
        ],
        "enhanced_features": [
            "Robust error handling and recovery",
            "Automatic data type detection and conversion",
            "Smart column matching and access",
            "LLM-based code generation and fixing",
            "Hybrid OCR approach with multiple fallbacks",
            "Comprehensive file format support",
            "Intelligent data extraction from any source",
            "Question-aware response formatting"
        ]
    })

@app.errorhandler(404)
def not_found(error):
    return jsonify({
        "error": "Endpoint not found",
        "available_endpoints": ["/api/", "/health", "/capabilities"],
        "help": "Use POST /api/ to analyze data, GET /health for status, GET /capabilities for supported formats"
    }), 404

@app.errorhandler(500)
def internal_error(error):
    return jsonify({
        "error": "Internal server error",
        "message": "The server encountered an unexpected condition",
        "help": "Please check your request format and data files"
    }), 500

if __name__ == '__main__':
    # Load environment variables
    load_dotenv()
    
    # Validate required environment variables
    if not os.getenv('ANTHROPIC_API_KEY'):
        logger.error("ANTHROPIC_API_KEY environment variable is required")
        exit(1)
    
    # Configure server
    port = int(os.getenv('PORT', 5000))
    debug = os.getenv('FLASK_DEBUG', 'False').lower() == 'true'
    host = os.getenv('FLASK_HOST', '0.0.0.0')
    
    logger.info(f"Starting Universal Data Analyst Agent v5.0 (Comprehensive) on {host}:{port}")
    logger.info(f"Debug mode: {debug}")
    logger.info("Key improvements: Comprehensive file support, hybrid OCR, enhanced error handling")
    logger.info("Supported formats: CSV, Excel, JSON, Images (PNG/JPG/etc), PDFs, Word, PowerPoint, SQLite, ZIP archives, and more")
    
    # Run the Flask app
    app.run(
        host=host,
        port=port,
        debug=debug,
        threaded=True,
        use_reloader=False
    )